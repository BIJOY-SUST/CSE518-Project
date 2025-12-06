"""
From Gestures to Words: A Hybrid System for Interactive PowerPoint Presentations
- Navigate slides with hand gestures
- Draw on slides with index finger, erase with 3 fingers
- Laser pointer mode with peace sign gesture
- Real-time speech-to-text display
- No PowerPoint installation required
"""
import sys
import cv2
import numpy as np
import mediapipe as mp
import os
import time
import queue
import threading
import subprocess

from PIL import Image
from collections import deque
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
from pdf2image import convert_from_path


# Speech Recognition
try:
    import speech_recognition as sr
    SPEECH_AVAILABLE = True
except ImportError:
    SPEECH_AVAILABLE = False
    print("speech_recognition not installed. Speech-to-text will be disabled.")
    print("Install with: pip install SpeechRecognition pyaudio")

# ─────────────────────────────────────────────────────────────────────────────
# CONFIGURATION
# ─────────────────────────────────────────────────────────────────────────────

WINDOW_WIDTH = 1280
WINDOW_HEIGHT = 720
GESTURE_COOLDOWN = 0.5
DRAWING_THICKNESS = 3
ERASER_SIZE = 30

# Available drawing colors (BGR format)
DRAWING_COLORS = {
    'red': (0, 0, 255),
    'green': (0, 255, 0),
    'blue': (255, 0, 0),
    'yellow': (0, 255, 255),
    'cyan': (255, 255, 0),
    'magenta': (255, 0, 255),
    'white': (255, 255, 255),
    'orange': (0, 165, 255),
    'purple': (255, 0, 127),
    'black': (0, 0, 0)
}

COLOR_NAMES = list(DRAWING_COLORS.keys())
DEFAULT_COLOR_INDEX = 0  # Start with red

# ─────────────────────────────────────────────────────────────────────────────
# MEDIAPIPE SETUP
# ─────────────────────────────────────────────────────────────────────────────

mp_hands = mp.solutions.hands
mp_drawing = mp.solutions.drawing_utils
mp_drawing_styles = mp.solutions.drawing_styles

hands = mp_hands.Hands(
    static_image_mode=False,
    max_num_hands=1,
    min_detection_confidence=0.7,
    min_tracking_confidence=0.7,
    model_complexity=1
)

# ─────────────────────────────────────────────────────────────────────────────
# PPTX TO IMAGE CONVERTER
# ─────────────────────────────────────────────────────────────────────────────

def pptx_to_images_pdf(pptx_path):
    """
    Convert PPTX to PDF first, then to images
    """
    try:
        # Convert PPTX to PDF using LibreOffice
        pdf_path = pptx_path.replace('.pptx', '_temp.pdf')
        pptx_dir = os.path.dirname(pptx_path)
        
        print("🔄 Converting PPTX to PDF for better quality...")
        
        # Try LibreOffice - check macOS location first
        libreoffice_paths = [
            '/Applications/LibreOffice.app/Contents/MacOS/soffice',  # macOS
            'libreoffice',  # Linux/Windows or if in PATH
            'soffice'  # Alternative command
        ]
        
        libreoffice_cmd = None
        for path in libreoffice_paths:
            if os.path.exists(path) or path in ['libreoffice', 'soffice']:
                libreoffice_cmd = path
                break
        
        if not libreoffice_cmd:
            print("⚠️  LibreOffice not found")
            return None
        
        # Try LibreOffice conversion
        result = subprocess.run([
            libreoffice_cmd, '--headless', '--convert-to', 'pdf',
            '--outdir', pptx_dir, pptx_path
        ], capture_output=True, text=True, timeout=60)
        
        # Check if PDF was created
        possible_pdf = pptx_path.replace('.pptx', '.pdf')
        if os.path.exists(possible_pdf):
            pdf_path = possible_pdf
        
        if os.path.exists(pdf_path):
            print("✅ PDF created successfully")
            print("🔄 Converting PDF to images...")
            
            # Convert PDF to images
            images = convert_from_path(pdf_path, dpi=150)
            slides = []
            
            for idx, image in enumerate(images):
                print(f"  Processing slide {idx + 1}...", end='\r')
                # Resize to window dimensions
                img = image.resize((WINDOW_WIDTH, WINDOW_HEIGHT), Image.Resampling.LANCZOS)
                img_array = np.array(img)
                img_bgr = cv2.cvtColor(img_array, cv2.COLOR_RGB2BGR)
                slides.append(img_bgr)
            
            # Clean up temporary PDF
            try:
                if pdf_path != possible_pdf:
                    os.remove(pdf_path)
            except:
                pass
            print(f"\n✅ Loaded {len(slides)} high-quality slides!")
            return slides
        else:
            print("⚠️  PDF conversion failed")
            return None
            
    except ImportError:
        print("⚠️  pdf2image not installed")
        return None
    except subprocess.TimeoutExpired:
        print("⚠️  PDF conversion timeout")
        return None
    except Exception as e:
        print(f"⚠️  PDF conversion error: {e}")
        return None

def load_pptx(pptx_path):
    """Load PPTX file"""
    slides = pptx_to_images_pdf(pptx_path)
    
    return slides

# ─────────────────────────────────────────────────────────────────────────────
# GESTURE RECOGNITION & DISTANCE ESTIMATION
# ─────────────────────────────────────────────────────────────────────────────

def estimate_hand_distance(hand_landmarks, frame_width, frame_height):
    """
    Estimate distance from hand to camera in centimeters
    Uses hand size (width) and z-coordinate from MediaPipe
    
    Returns: (distance_cm, hand_width_pixels)
    """
    # Calculate hand bounding box in pixels
    x_coords = [lm.x for lm in hand_landmarks.landmark]
    y_coords = [lm.y for lm in hand_landmarks.landmark]
    
    x_min = min(x_coords) * frame_width
    x_max = max(x_coords) * frame_width
    y_min = min(y_coords) * frame_height
    y_max = max(y_coords) * frame_height
    
    hand_width_px = x_max - x_min
    hand_height_px = y_max - y_min
    
    # Average adult hand width is approximately 8-10 cm
    # We'll use 9 cm as a reasonable average
    REAL_HAND_WIDTH_CM = 9.0
    
    # Focal length estimation (calibrated for typical webcam at 640x480)
    # This can be adjusted based on your specific camera
    # For better accuracy: focal_length ≈ (pixel_width * known_distance) / real_width
    FOCAL_LENGTH = 600  # Typical for 640x480 webcam
    
    # Distance formula: distance = (real_width * focal_length) / pixel_width
    if hand_width_px > 0:
        distance_cm = (REAL_HAND_WIDTH_CM * FOCAL_LENGTH) / hand_width_px
        
        # Use z-coordinate for additional refinement
        # MediaPipe's z is relative depth (negative means closer, positive means farther)
        # Z typically ranges from -0.15 (very close) to 0.15 (far)
        avg_z = sum([lm.z for lm in hand_landmarks.landmark]) / len(hand_landmarks.landmark)
        
        # Apply z-based adjustment (more conservative than before)
        # When z is negative (closer), we slightly reduce distance
        # When z is positive (farther), we slightly increase distance
        z_factor = 1.0 + (avg_z * 0.5)  # More conservative multiplier
        z_factor = max(0.7, min(1.3, z_factor))  # Clamp z_factor to reasonable range
        
        distance_cm *= z_factor
        
        # Clamp to reasonable range (5 cm to 250 cm) - Extended range
        distance_cm = max(5, min(250, distance_cm))
        
        return distance_cm, hand_width_px
    
    return None, 0

def count_fingers(hand_landmarks):
    """Count extended fingers"""
    finger_tips = [8, 12, 16, 20]  # Index, Middle, Ring, Pinky
    thumb_tip = 4
    
    fingers = []
    
    # Thumb (check x-coordinate for left/right hand)
    if hand_landmarks.landmark[thumb_tip].x < hand_landmarks.landmark[thumb_tip - 1].x:
        fingers.append(1)
    else:
        fingers.append(0)
    
    # Other fingers (check y-coordinate)
    for tip in finger_tips:
        if hand_landmarks.landmark[tip].y < hand_landmarks.landmark[tip - 2].y:
            fingers.append(1)
        else:
            fingers.append(0)
    
    return tuple(fingers)

def detect_gesture(hand_landmarks):
    """Detect specific gestures"""
    fingers = count_fingers(hand_landmarks)
    
    # Gesture mappings
    if fingers == (0, 1, 0, 0, 0):  # Index finger only
        return 'draw'
    elif fingers == (0, 1, 1, 0, 0):  # Index + Middle (peace sign)
        return 'laser_pointer'
    elif fingers == (0, 1, 1, 1, 0):  # Index + Middle + Ring (erase)
        return 'erase'
    elif fingers == (1, 0, 0, 0, 0):  # Thumb only
        return 'previous'
    elif fingers == (1, 1, 1, 1, 1):  # All fingers (palm)
        return 'next'
    elif fingers == (0, 0, 0, 0, 0):  # Fist - do nothing
        return None
    elif fingers == (1, 1, 0, 0, 1):  # Thumb + Index + Pinky
        return 'clear_all'
    elif fingers == (0, 1, 0, 0, 1):  # Index + Pinky
        return 'change_color'
    
    return None

# ─────────────────────────────────────────────────────────────────────────────
# SPEECH RECOGNITION CLASS
# ─────────────────────────────────────────────────────────────────────────────

class SpeechTranscriber:
    """Real-time speech-to-text transcription in background thread"""
    
    def __init__(self):
        self.transcript_queue = queue.Queue()
        self.is_running = False
        self.thread = None
        self.current_transcript = ""
        self.transcript_history = deque(maxlen=10)  # Keep last 10 phrases
        
    def start(self):
        """Start background speech recognition thread"""
        if not SPEECH_AVAILABLE:
            print("⚠️  Speech recognition not available")
            return False
        
        self.is_running = True
        self.thread = threading.Thread(target=self._listen_loop, daemon=True)
        self.thread.start()
        print("🎤 Speech recognition started")
        return True
    
    def stop(self):
        """Stop background speech recognition"""
        self.is_running = False
        if self.thread:
            self.thread.join(timeout=2)
        print("🎤 Speech recognition stopped")
    
    def _listen_loop(self):
        """Background thread listening for speech"""
        recognizer = sr.Recognizer()
        recognizer.energy_threshold = 4000  # Adjust for ambient noise
        recognizer.dynamic_energy_threshold = True
        recognizer.pause_threshold = 0.8  # Seconds of silence to consider end of phrase
        
        with sr.Microphone() as source:
            print("🎤 Adjusting for ambient noise... Please wait.")
            recognizer.adjust_for_ambient_noise(source, duration=2)
            print("🎤 Ready! Start speaking...")
            
            while self.is_running:
                try:
                    # Listen for audio
                    audio = recognizer.listen(source, timeout=1, phrase_time_limit=10)
                    
                    # Recognize speech using Google Speech Recognition
                    try:
                        text = recognizer.recognize_google(audio)
                        if text:
                            self.transcript_queue.put(text)
                            print(f"🎤 Recognized: {text}")
                    except sr.UnknownValueError:
                        # Could not understand audio
                        pass
                    except sr.RequestError as e:
                        print(f"⚠️  Speech recognition error: {e}")
                        
                except sr.WaitTimeoutError:
                    # No speech detected, continue listening
                    continue
                except Exception as e:
                    if self.is_running:
                        print(f"⚠️  Speech error: {e}")
    
    def get_transcript(self):
        """Get latest transcript text"""
        # Update transcript with any new phrases
        while not self.transcript_queue.empty():
            try:
                text = self.transcript_queue.get_nowait()
                self.transcript_history.append(text)
                self.current_transcript = text
            except queue.Empty:
                break
        
        return self.current_transcript
    
    def get_full_history(self):
        """Get all recent transcript phrases"""
        return list(self.transcript_history)

# ─────────────────────────────────────────────────────────────────────────────
# PRESENTATION VIEWER CLASS
# ─────────────────────────────────────────────────────────────────────────────

class PresentationViewer:
    def __init__(self, slides):
        self.slides = slides
        self.current_slide = 0
        self.drawing_layers = [np.zeros((WINDOW_HEIGHT, WINDOW_WIDTH, 3), dtype=np.uint8) 
                               for _ in slides]
        self.last_action_time = 0
        self.drawing_mode = False
        self.last_draw_point = None
        self.pointer_trail = deque(maxlen=10)
        self.laser_pointer_pos = None
        self.current_color_index = DEFAULT_COLOR_INDEX
        self.current_color = DRAWING_COLORS[COLOR_NAMES[self.current_color_index]]
    
    def get_current_color_name(self):
        """Get the name of current drawing color"""
        return COLOR_NAMES[self.current_color_index]
    
    def cycle_color(self):
        """Cycle to next drawing color"""
        self.current_color_index = (self.current_color_index + 1) % len(COLOR_NAMES)
        self.current_color = DRAWING_COLORS[COLOR_NAMES[self.current_color_index]]
        print(f"Drawing color: {self.get_current_color_name().upper()}")
        return self.get_current_color_name()
        
    def get_current_frame(self, show_laser=False, show_eraser_cursor=False, eraser_pos=None, 
                          transcript_text="", transcript_history=None):
        """Get current slide with drawings overlay, laser pointer, eraser cursor, and transcript"""
        slide = self.slides[self.current_slide].copy()
        drawing = self.drawing_layers[self.current_slide]
        
        # Overlay drawing
        mask = cv2.cvtColor(drawing, cv2.COLOR_BGR2GRAY)
        _, mask = cv2.threshold(mask, 1, 255, cv2.THRESH_BINARY)
        mask_inv = cv2.bitwise_not(mask)
        
        slide_bg = cv2.bitwise_and(slide, slide, mask=mask_inv)
        drawing_fg = cv2.bitwise_and(drawing, drawing, mask=mask)
        
        result = cv2.add(slide_bg, drawing_fg)
        
        # Draw laser pointer if active
        if show_laser and self.laser_pointer_pos is not None:
            x, y = self.laser_pointer_pos
            # Draw red dot with glow effect (smaller size)
            cv2.circle(result, (x, y), 15, (0, 0, 255), -1)  # Outer glow
            cv2.circle(result, (x, y), 12, (0, 0, 220), -1)  # Middle
            cv2.circle(result, (x, y), 9, (0, 0, 180), -1)   # Inner
            cv2.circle(result, (x, y), 6, (0, 0, 255), -1)   # Center bright
        
        # Draw eraser cursor if active
        if show_eraser_cursor and eraser_pos is not None:
            x, y = eraser_pos
            # Draw semi-transparent orange circle showing eraser area
            overlay = result.copy()
            # Outer circle (eraser boundary)
            cv2.circle(overlay, (x, y), ERASER_SIZE, (0, 165, 255), 2)  # Orange outline
            # Inner filled circle (semi-transparent)
            cv2.circle(overlay, (x, y), ERASER_SIZE, (0, 165, 255), -1)  # Orange fill
            # Blend with transparency
            cv2.addWeighted(overlay, 0.3, result, 0.7, 0, result)
            # Draw crosshair for precision
            cv2.line(result, (x - 10, y), (x + 10, y), (0, 165, 255), 2)
            cv2.line(result, (x, y - 10), (x, y + 10), (0, 165, 255), 2)
        
        # Add transcript bar BELOW slide content
        if SPEECH_AVAILABLE and (transcript_text or transcript_history):
            result = self._add_transcript_bar(result, transcript_text, transcript_history)
        
        return result
    
    def _add_transcript_bar(self, slide_frame, transcript_text, transcript_history):
        """Add YouTube-style subtitle overlay at bottom of slide"""
        slide_height, slide_width = slide_frame.shape[:2]
        result = slide_frame.copy()
        
        # Get the text to display
        display_text = ""
        if transcript_text:
            display_text = transcript_text
        elif transcript_history:
            display_text = transcript_history[-1] if transcript_history else ""
        
        if not display_text:
            return result
        
        # Split text into lines if too long (max ~50 chars per line for readability)
        max_chars_per_line = 50
        words = display_text.split()
        lines = []
        current_line = ""
        
        for word in words:
            test_line = current_line + " " + word if current_line else word
            if len(test_line) <= max_chars_per_line:
                current_line = test_line
            else:
                if current_line:
                    lines.append(current_line)
                current_line = word
        if current_line:
            lines.append(current_line)
        
        # Limit to last 2 lines only
        if len(lines) > 2:
            lines = lines[-2:]
        
        # Font settings
        font = cv2.FONT_HERSHEY_SIMPLEX
        font_scale = 0.9
        font_thickness = 2
        line_height = 40
        padding_x = 15
        padding_y = 8
        
        # Calculate total background dimensions
        max_text_width = 0
        for line in lines:
            (text_w, text_h), baseline = cv2.getTextSize(line, font, font_scale, font_thickness)
            max_text_width = max(max_text_width, text_w)
        
        total_text_height = len(lines) * line_height
        
        # Position at bottom center
        y_start = slide_height - 30 - total_text_height
        
        # Calculate background rectangle (covers all lines)
        bg_x1 = (slide_width - max_text_width) // 2 - padding_x
        bg_y1 = y_start - padding_y
        bg_x2 = (slide_width + max_text_width) // 2 + padding_x
        bg_y2 = y_start + total_text_height + padding_y
        
        # Draw semi-transparent black background (slide visible through it)
        overlay = result.copy()
        cv2.rectangle(overlay, (bg_x1, bg_y1), (bg_x2, bg_y2), (0, 0, 0), -1)
        cv2.addWeighted(overlay, 0.5, result, 0.5, 0, result)  # 50% transparent
        
        # Draw each line of white text centered
        for i, line in enumerate(lines):
            (text_w, text_h), baseline = cv2.getTextSize(line, font, font_scale, font_thickness)
            text_x = (slide_width - text_w) // 2
            text_y = y_start + (i + 1) * line_height - 10
            
            # Draw white text
            cv2.putText(result, line, (text_x, text_y), font, font_scale, (255, 255, 255), font_thickness)
        
        return result
    
    def next_slide(self):
        if time.time() - self.last_action_time > GESTURE_COOLDOWN:
            if self.current_slide < len(self.slides) - 1:
                self.current_slide += 1
                self.last_action_time = time.time()
                print(f"➡️  Slide {self.current_slide + 1}/{len(self.slides)}")
                return True
        return False
    
    def previous_slide(self):
        if time.time() - self.last_action_time > GESTURE_COOLDOWN:
            if self.current_slide > 0:
                self.current_slide -= 1
                self.last_action_time = time.time()
                print(f"⬅️  Slide {self.current_slide + 1}/{len(self.slides)}")
                return True
        return False
    
    def update_laser_pointer(self, x, y):
        """Update laser pointer position"""
        self.laser_pointer_pos = (x, y)
    
    def clear_laser_pointer(self):
        """Clear laser pointer"""
        self.laser_pointer_pos = None
    
    def draw_point(self, x, y):
        """Draw on current slide using current color"""
        drawing = self.drawing_layers[self.current_slide]
        
        if self.last_draw_point is not None:
            cv2.line(drawing, self.last_draw_point, (x, y), 
                    self.current_color, DRAWING_THICKNESS)
        
        self.last_draw_point = (x, y)
    
    def erase_area(self, x, y):
        """Erase drawing in area"""
        drawing = self.drawing_layers[self.current_slide]
        cv2.circle(drawing, (x, y), ERASER_SIZE, (0, 0, 0), -1)
    
    def clear_all_drawings(self):
        """Clear all drawings on current slide"""
        self.drawing_layers[self.current_slide] = np.zeros(
            (WINDOW_HEIGHT, WINDOW_WIDTH, 3), dtype=np.uint8)
        print("🗑️  Cleared all drawings on current slide")

# ─────────────────────────────────────────────────────────────────────────────
# UI OVERLAY FUNCTIONS
# ─────────────────────────────────────────────────────────────────────────────

def draw_gesture_info(frame, gesture, slide_info, fingers=None, current_color_name=None, distance_cm=None):
    """Draw gesture information on camera feed"""
    height, width = frame.shape[:2]
    
    # Display fingers array at top
    if fingers is not None:
        fingers_text = f"Fingers: {{{fingers[0]}, {fingers[1]}, {fingers[2]}, {fingers[3]}, {fingers[4]}}}"
        cv2.putText(frame, fingers_text, (10, 30), 
                   cv2.FONT_HERSHEY_SIMPLEX, 0.7, (255, 255, 255), 2)
    else:
        cv2.putText(frame, "Fingers: {0, 0, 0, 0, 0}", (10, 30), 
                   cv2.FONT_HERSHEY_SIMPLEX, 0.7, (150, 150, 150), 2)
    
    # Display distance BELOW MIC status (moved from top)
    if distance_cm is not None:
        distance_text = f"Distance: {distance_cm:.1f} cm"
        # Color code: green (30-60cm optimal), yellow (20-80cm), red (otherwise)
        if 30 <= distance_cm <= 60:
            distance_color = (0, 255, 0)  # Green - optimal
        elif 20 <= distance_cm <= 80:
            distance_color = (0, 255, 255)  # Yellow - acceptable
        else:
            distance_color = (0, 0, 255)  # Red - too far/close
        
        # Position below MIC status (moved left and down)
        box_x1 = width - 250
        box_y1 = 50  # Below MIC: ON (which is at y=30)
        box_x2 = width - 10
        box_y2 = 90
        
        # Draw distance with background for better visibility
        cv2.rectangle(frame, (box_x1, box_y1), (box_x2, box_y2), (0, 0, 0), -1)
        cv2.rectangle(frame, (box_x1, box_y1), (box_x2, box_y2), distance_color, 2)
        cv2.putText(frame, distance_text, (box_x1 + 10, box_y1 + 28), 
                   cv2.FONT_HERSHEY_SIMPLEX, 0.7, distance_color, 2)
    
    # Gesture indicator with specific patterns
    gesture_patterns = {
        'draw': '{0,1,0,0,0}',
        'laser_pointer': '{0,1,1,0,0}',
        'erase': '{0,1,1,1,0}',
        'previous': '{1,0,0,0,0}',
        'next': '{1,1,1,1,1}',
        'change_color': '{0,1,0,0,1}',
        'clear_all': '{1,1,0,0,1}'
    }
    
    if gesture and gesture in gesture_patterns:
        gesture_text = gesture.upper().replace('_', ' ')
        pattern = gesture_patterns[gesture]
        if gesture == 'draw':
            color = (0, 255, 0)  # Green
        elif gesture == 'laser_pointer':
            color = (0, 0, 255)  # Red
        elif gesture == 'erase':
            color = (0, 165, 255)  # Orange
        elif gesture == 'change_color':
            color = (255, 0, 255)  # Magenta
        else:
            color = (255, 255, 0)  # Yellow
        # Removed emoji - just show gesture name
        cv2.putText(frame, f"Mode: {gesture_text}", (10, 65),
                   cv2.FONT_HERSHEY_SIMPLEX, 0.8, color, 2)
    
    # Slide info and color indicator
    slide_text = slide_info
    if current_color_name:
        slide_text += f" | Color: {current_color_name.upper()}"
    cv2.putText(frame, slide_text, (10, 95),
               cv2.FONT_HERSHEY_SIMPLEX, 0.6, (255, 255, 255), 2)
    
    # Instructions at bottom with gesture patterns
    instructions = [
        "{0,1,0,0,0}=Draw | {0,1,1,0,0}=Laser | {0,1,1,1,0}=Erase | {0,1,0,0,1}=Color",
        "{1,0,0,0,0}=Prev | {1,1,1,1,1}=Next | {1,1,0,0,1}=Clear | Q=Quit",
        "N=Next Slide | P=Previous Slide | C=Clear | Distance: Top Right"
    ]
    
    y_pos = height - 70
    for instruction in instructions:
        cv2.putText(frame, instruction, (10, y_pos), 
                   cv2.FONT_HERSHEY_SIMPLEX, 0.45, (255, 255, 255), 1)
        y_pos += 30

def draw_speech_status(frame):
    """Draw simple speech recognition status on camera window"""
    height, width = frame.shape[:2]
    
    # Show speech status if available (top right corner)
    if SPEECH_AVAILABLE:
        mic_status = "MIC: ON"
        mic_color = (0, 255, 0)  # Green
        cv2.putText(frame, mic_status, (width - 120, 30),
                   cv2.FONT_HERSHEY_SIMPLEX, 0.6, mic_color, 2)

# ─────────────────────────────────────────────────────────────────────────────
# MAIN APPLICATION
# ─────────────────────────────────────────────────────────────────────────────

def main():
    print("\n" + "="*80)
    print("🎯 LOCAL PPTX HAND GESTURE PRESENTATION VIEWER")
    print("="*80)
    
    # File selection
    print("\n📂 Enter the path to your PPTX file:")
    print("   (You can drag and drop the file here)")
    print("   Example: /Users/bijoy/Documents/presentation.pptx")
    pptx_path = input("\nPath: ").strip().strip('"').strip("'")
    
    # Validate file
    if not pptx_path:
        print("❌ No file path provided")
        return
    
    if not os.path.exists(pptx_path):
        print(f"❌ File not found: {pptx_path}")
        return
    
    if not pptx_path.lower().endswith('.pptx'):
        print("❌ Please provide a .pptx file")
        return
    
    # Load presentation
    print("\n🔄 Loading presentation...")
    slides = load_pptx(pptx_path)
    
    if slides is None or len(slides) == 0:
        print("❌ Failed to load presentation")
        return
    
    # Initialize viewer
    viewer = PresentationViewer(slides)
    
    # Open camera
    print("\n📷 Initializing camera...")
    cap = cv2.VideoCapture(0)
    cap.set(cv2.CAP_PROP_FRAME_WIDTH, 640)
    cap.set(cv2.CAP_PROP_FRAME_HEIGHT, 480)
    
    if not cap.isOpened():
        print("❌ Cannot access camera")
        return
    
    print("✅ Camera initialized")
    
    # Initialize speech transcriber
    transcriber = None
    if SPEECH_AVAILABLE:
        print("\n🎤 Initializing speech recognition...")
        transcriber = SpeechTranscriber()
        if transcriber.start():
            print("✅ Speech-to-text ready")
        else:
            transcriber = None
    else:
        print("\n⚠️  Speech recognition disabled (install: pip install SpeechRecognition pyaudio)")
    
    print("\n" + "="*80)
    print("✋ HAND GESTURES:")
    print("="*80)
    print("  Index finger only     → Draw on slide")
    print("  Index + Middle        → Laser pointer (red dot)")
    print("  Index+Middle+Ring     → Erase (move to erase area)")
    print("  Thumb only            → Previous slide")
    print("  All fingers (palm)    → Next slide")
    print("  Thumb+Index+Pinky     → Clear all drawings")
    print("\n⌨️  KEYBOARD:")
    print("  'q' → Quit")
    print("  'n' → Next slide")
    print("  'p' → Previous slide")
    print("  'c' → Clear all drawings on current slide")
    print("\n🎨 DRAWING COLORS:")
    print("  Red, Green, Blue, Yellow, Cyan, Magenta, White, Orange, Purple, Black")
    print("  Make {0,1,0,0,1} gesture (Index + Pinky) to cycle through colors")
    print("\n📏 DISTANCE MEASUREMENT:")
    print("  Real-time distance displayed in top-right corner of camera window")
    print("  🟢 Green (30-60cm): Optimal distance for gesture recognition")
    print("  🟡 Yellow (20-80cm): Acceptable range")
    print("  🔴 Red: Too close or too far - adjust your position")
    if SPEECH_AVAILABLE:
        print("\n🎤 SPEECH-TO-TEXT:")
        print("  Real-time transcription displayed in PPTX window")
    print("="*80 + "\n")
    
    print("▶️  Starting presentation...")
    print(f"📊 Showing slide 1/{len(slides)}\n")

    # Create windows - slide + transcript bar
    cv2.namedWindow('Presentation', cv2.WINDOW_NORMAL)
    cv2.namedWindow('Camera (Hand Tracking)', cv2.WINDOW_NORMAL)
    
    last_gesture = None
    gesture_start_time = time.time()
    
    try:
        while True:
            # Read camera frame
            ret, frame = cap.read()
            if not ret:
                print("❌ Failed to read camera frame")
                break
            
            frame = cv2.flip(frame, 1)
            rgb_frame = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
            
            # Process hand landmarks
            results = hands.process(rgb_frame)
            
            current_gesture = None
            current_fingers = None
            eraser_position = None
            current_distance = None
            
            if results.multi_hand_landmarks:
                for hand_landmarks in results.multi_hand_landmarks:
                    # Draw hand landmarks
                    mp_drawing.draw_landmarks(
                        frame,
                        hand_landmarks,
                        mp_hands.HAND_CONNECTIONS,
                        mp_drawing_styles.get_default_hand_landmarks_style(),
                        mp_drawing_styles.get_default_hand_connections_style()
                    )
                    
                    # Calculate hand distance from camera
                    frame_h, frame_w = frame.shape[:2]
                    current_distance, _ = estimate_hand_distance(hand_landmarks, frame_w, frame_h)
                    
                    # Get fingers array
                    current_fingers = count_fingers(hand_landmarks)
                    
                    # Detect gesture
                    current_gesture = detect_gesture(hand_landmarks)
                    
                    # Get index finger tip position for drawing/pointing
                    index_tip = hand_landmarks.landmark[8]
                    x = int(index_tip.x * WINDOW_WIDTH)
                    y = int(index_tip.y * WINDOW_HEIGHT)
                    
                    # Handle gestures
                    if current_gesture == 'draw':
                        viewer.draw_point(x, y)
                        viewer.clear_laser_pointer()
                    elif current_gesture == 'laser_pointer':
                        viewer.update_laser_pointer(x, y)
                        viewer.last_draw_point = None
                    elif current_gesture == 'erase':
                        viewer.erase_area(x, y)
                        eraser_position = (x, y)  # Track position for cursor
                        viewer.last_draw_point = None
                        viewer.clear_laser_pointer()
                    elif current_gesture == 'change_color':
                        if last_gesture != 'change_color' or time.time() - gesture_start_time > 1:
                            viewer.cycle_color()
                            gesture_start_time = time.time()
                        viewer.last_draw_point = None
                        viewer.clear_laser_pointer()
                    elif current_gesture == 'next':
                        if last_gesture != 'next' or time.time() - gesture_start_time > 1:
                            viewer.next_slide()
                            gesture_start_time = time.time()
                        viewer.last_draw_point = None
                        viewer.clear_laser_pointer()
                    elif current_gesture == 'previous':
                        if last_gesture != 'previous' or time.time() - gesture_start_time > 1:
                            viewer.previous_slide()
                            gesture_start_time = time.time()
                        viewer.last_draw_point = None
                        viewer.clear_laser_pointer()
                    elif current_gesture == 'clear_all':
                        if last_gesture != 'clear_all' or time.time() - gesture_start_time > 1:
                            viewer.clear_all_drawings()
                            gesture_start_time = time.time()
                        viewer.last_draw_point = None
                        viewer.clear_laser_pointer()
                    else:
                        viewer.last_draw_point = None
                        viewer.clear_laser_pointer()
            else:
                viewer.last_draw_point = None
                viewer.clear_laser_pointer()
            
            last_gesture = current_gesture
            
            # Get speech transcript if available
            current_transcript = ""
            transcript_history = []
            if transcriber:
                current_transcript = transcriber.get_transcript()
                transcript_history = transcriber.get_full_history()
            
            # Draw UI overlay on camera feed
            slide_info = f"Slide {viewer.current_slide + 1}/{len(slides)}"
            draw_gesture_info(frame, current_gesture, slide_info, current_fingers, 
                            viewer.get_current_color_name(), current_distance)
            
            # Draw speech status indicator on camera window
            draw_speech_status(frame)
            
            # Get current presentation frame with visual indicators AND transcript
            show_laser = (current_gesture == 'laser_pointer')
            show_eraser = (current_gesture == 'erase')
            presentation_frame = viewer.get_current_frame(
                show_laser=show_laser,
                show_eraser_cursor=show_eraser,
                eraser_pos=eraser_position,
                transcript_text=current_transcript,
                transcript_history=transcript_history
            )
            
            # Show windows
            cv2.imshow('Camera (Hand Tracking)', frame)
            cv2.imshow('Presentation', presentation_frame)
            
            # Handle keyboard input
            key = cv2.waitKey(1) & 0xFF
            if key == ord('q'):
                print("\n👋 Quitting...")
                break
            elif key == ord('n'):
                viewer.next_slide()
            elif key == ord('p'):
                viewer.previous_slide()
            elif key == ord('c'):
                viewer.clear_all_drawings()
    
    except KeyboardInterrupt:
        print("\n\n⚠️  Interrupted by user")
    
    finally:
        # Cleanup
        if transcriber:
            transcriber.stop()
        cap.release()
        cv2.destroyAllWindows()
        hands.close()
        print("\n✅ Application closed. Thank you!")

if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        print(f"\n❌ Error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
